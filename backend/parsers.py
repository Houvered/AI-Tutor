"""Document parsers.

Each parser reads one file type and returns plain text. A single `parse_file`
function picks the right parser based on the file extension.
"""

import os

from backend.config import ALLOWED_EXTENSIONS


def _parse_pdf(path: str) -> str:
    import fitz  # PyMuPDF

    doc = fitz.open(path)
    pages = []
    for page in doc:
        pages.append(page.get_text())
    doc.close()
    return "\n\n".join(pages)


def _parse_docx(path: str) -> str:
    import docx

    document = docx.Document(path)
    parts = [p.text for p in document.paragraphs if p.text.strip()]
    return "\n\n".join(parts)


def _parse_pptx(path: str) -> str:
    from pptx import Presentation

    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        texts = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
        if texts:
            parts.append("\n".join(texts))
    return "\n\n".join(parts)


def _parse_text(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


_PARSERS = {
    ".pdf": _parse_pdf,
    ".docx": _parse_docx,
    ".pptx": _parse_pptx,
    ".txt": _parse_text,
    ".md": _parse_text,
}


def parse_file(path: str) -> str:
    """Return the plain text of a file, based on its extension."""
    ext = os.path.splitext(path)[1].lower()
    parser = _PARSERS.get(ext)
    if parser is None:
        raise ValueError(
            f"Unsupported file type '{ext}'. Supported: {sorted(ALLOWED_EXTENSIONS)}"
        )
    return parser(path)